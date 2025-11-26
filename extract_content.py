import os
from docx import Document
from pptx import Presentation

def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading {docx_path}: {str(e)}"

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading {pptx_path}: {str(e)}"

def main():
    files = [f for f in os.listdir('.') if f.endswith(('.docx', '.pptx'))]
    
    with open('project_content_extracted.txt', 'w', encoding='utf-8') as f:
        for file in files:
            f.write(f"--- START OF {file} ---\n")
            if file.endswith('.docx'):
                content = extract_text_from_docx(file)
            elif file.endswith('.pptx'):
                content = extract_text_from_pptx(file)
            f.write(content)
            f.write(f"\n--- END OF {file} ---\n\n")

if __name__ == "__main__":
    main()

